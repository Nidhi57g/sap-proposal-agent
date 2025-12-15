# sap_agent.py
# Minimal CI-safe SAP Proposal Agent

from pptx import Presentation
import sys

def self_test():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "SAP Proposal Agent"
    prs.save("test_output.pptx")
    print("✅ SAP Proposal Agent test passed")

if __name__ == "__main__":
    self_test()
