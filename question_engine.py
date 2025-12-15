from pptx import Presentation
import os
from collections import Counter

KNOWLEDGE_DIR = "knowledge"

COMMON_SECTIONS = [
    "Scope",
    "Business Benefits",
    "Implementation Approach",
    "Timeline",
    "Assumptions",
    "Out of Scope",
    "Commercials",
    "Pricing",
    "Risks",
    "Dependencies"
]

def extract_slide_titles():
    titles = []

    for file in os.listdir(KNOWLEDGE_DIR):
        if file.endswith(".pptx"):
            prs = Presentation(os.path.join(KNOWLEDGE_DIR, file))
            for slide in prs.slides:
                if slide.shapes.title:
                    titles.append(slide.shapes.title.text.strip())

    return titles

def generate_questions():
    titles = extract_slide_titles()
    counter = Counter(titles)

    questions = []

    for section in COMMON_SECTIONS:
        if not any(section.lower() in t.lower() for t in counter):
            questions.append(
                f"Please confirm details for the section: {section}"
            )

    # Business-specific questions
    questions.extend([
        "What is the client industry?",
        "Which SAP modules are in scope?",
        "Is this Greenfield, Brownfield, or Selective carve-out?",
        "Expected project duration?",
        "Any compliance or regulatory requirements?",
        "Target go-live date?",
        "Preferred pricing model (Fixed / T&M)?"
    ])

    return questions

if __name__ == "__main__":
    qs = generate_questions()
    print("🧠 Agent Clarification Questions:\n")
    for q in qs:
        print(f"- {q}")
