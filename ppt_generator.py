from pptx import Presentation

def generate_ppt(slides, filename="generated_proposal.pptx"):
    prs = Presentation()

    for slide_data in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = slide_data["title"]

        body = slide.placeholders[1].text_frame
        body.clear()

        for line in slide_data["content"]:
            p = body.add_paragraph()
            p.text = line

    prs.save(filename)
    print(f"📊 Proposal generated: {filename}")
